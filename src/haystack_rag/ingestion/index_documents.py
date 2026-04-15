from __future__ import annotations

import argparse
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from haystack import Document
from haystack.components.embedders import OpenAIDocumentEmbedder
from haystack.document_stores.types import DuplicatePolicy
from haystack_integrations.components.embedders.fastembed import FastembedDocumentEmbedder
from openpyxl import load_workbook
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from haystack_rag.config import AppConfig, create_document_store, require_secret

TEXT_EXTENSIONS = {
    ".csv",
    ".html",
    ".htm",
    ".json",
    ".log",
    ".md",
    ".rst",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}

RU_MARKERS = {"ru", "rus", "russian", "рус", "русский"}
EN_MARKERS = {"en", "eng", "english", "англ", "english-language"}
CYRILLIC_RE = re.compile(r"[А-Яа-яЁё]")


@dataclass(frozen=True)
class SourceUnit:
    text: str
    meta: dict[str, object]


@dataclass(frozen=True)
class ChunkSpan:
    content: str
    start: int
    end: int


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Index source documents into Qdrant.")
    parser.add_argument(
        "--input-dir",
        default=os.getenv("INPUT_DIR", "data/input"),
        help="Directory with source documents.",
    )
    parser.add_argument(
        "--recreate-index",
        action="store_true",
        help="Drop and recreate the Qdrant collection before indexing.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    config = AppConfig.from_env()
    input_dir = Path(args.input_dir).resolve()
    if not input_dir.exists():
        raise SystemExit(f"Input directory does not exist: {input_dir}")

    documents = list(build_documents(input_dir=input_dir, config=config))
    if not documents:
        raise SystemExit("No indexable documents were found.")

    embedder = create_document_embedder(config)
    embedded_documents = embedder.run(documents=documents)["documents"]

    document_store = create_document_store(config=config, recreate_index=args.recreate_index)
    written = document_store.write_documents(embedded_documents, policy=DuplicatePolicy.OVERWRITE)
    print(f"Indexed {written} chunks into '{config.qdrant_index}' from {input_dir}")


def build_documents(
    input_dir: Path,
    config: AppConfig,
) -> Iterable[Document]:
    for path in sorted(input_dir.rglob("*")):
        if not path.is_file():
            continue
        if path.name.startswith("."):
            continue

        source_units = extract_source_units(path=path, pdf_extractor=config.pdf_extractor)
        unit_texts = [unit.text for unit in source_units if unit.text.strip()]
        if not unit_texts:
            continue

        combined_text, page_spans = combine_source_units(source_units)
        metadata = build_metadata(path=path, input_dir=input_dir, text=combined_text)
        extractor = infer_extractor(source_units)
        chunks = chunk_text_with_offsets(
            text=combined_text,
            chunk_size=config.chunk_size,
            overlap=config.chunk_overlap,
        )

        for chunk_index, chunk in enumerate(chunks):
            yield Document(
                content=chunk.content,
                meta={
                    **metadata,
                    "extractor": extractor,
                    **resolve_page_metadata(page_spans=page_spans, start=chunk.start, end=chunk.end),
                    "chunk_index": chunk_index,
                },
            )


def create_document_embedder(config: AppConfig) -> OpenAIDocumentEmbedder | FastembedDocumentEmbedder:
    if config.embedding_provider == "openai":
        return OpenAIDocumentEmbedder(
            api_key=require_secret(config.embedding_api_key, "EMBEDDING_API_KEY or OPENAI_API_KEY"),
            api_base_url=config.embedding_api_base_url,
            dimensions=config.embedding_dimensions,
            model=config.embedding_model,
            progress_bar=True,
        )

    return FastembedDocumentEmbedder(
        model=config.embedding_model,
        cache_dir=config.fastembed_cache_path,
        prefix="passage: " if "e5" in config.embedding_model.lower() else "",
        progress_bar=True,
    )


def extract_source_units(path: Path, pdf_extractor: str = "hybrid") -> list[SourceUnit]:
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
        return [SourceUnit(text=text, meta={})] if text else []

    try:
        if suffix == ".pdf":
            return extract_pdf_units(path=path, mode=pdf_extractor)
        if suffix == ".docx":
            text = extract_docx_text(path)
            return [SourceUnit(text=text, meta={})] if text else []
        if suffix == ".pptx":
            text = extract_pptx_text(path)
            return [SourceUnit(text=text, meta={})] if text else []
        if suffix == ".xlsx":
            text = extract_xlsx_text(path)
            return [SourceUnit(text=text, meta={})] if text else []
    except Exception as exc:
        print(f"Skipping {path}: parser failed ({exc})")
        return []

    print(f"Skipping {path}: unsupported file type")
    return []


def extract_pdf_units(path: Path, mode: str = "hybrid") -> list[SourceUnit]:
    normalized_mode = mode.strip().lower()
    if normalized_mode == "pypdf":
        return extract_pdf_units_with_pypdf(path)

    if normalized_mode == "docling":
        return extract_pdf_units_with_docling(path) or extract_pdf_units_with_pypdf(path)

    pypdf_units = extract_pdf_units_with_pypdf(path)
    if not should_retry_pdf_with_docling(pypdf_units):
        return pypdf_units

    docling_units = extract_pdf_units_with_docling(path)
    if docling_units:
        return docling_units
    return pypdf_units


def extract_pdf_units_with_pypdf(path: Path) -> list[SourceUnit]:
    reader = PdfReader(str(path))
    units: list[SourceUnit] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = clean_pdf_text(page.extract_text() or "")
        if not text:
            continue
        units.append(
            SourceUnit(
                text=text,
                meta={
                    "page_number": page_number,
                    "page_start": page_number,
                    "page_end": page_number,
                    "page_label": f"p.{page_number}",
                    "extractor": "pypdf",
                },
            )
        )
    return units


def extract_pdf_units_with_docling(path: Path) -> list[SourceUnit]:
    try:
        from docling.document_converter import DocumentConverter
    except Exception as exc:
        print(f"Docling unavailable for {path}: {exc}")
        return []

    try:
        converter = DocumentConverter()
        result = converter.convert(str(path))
        text = clean_docling_text(result.document.export_to_markdown())
    except Exception as exc:
        print(f"Docling parser failed for {path}: {exc}")
        return []

    if not text:
        return []

    return [SourceUnit(text=text, meta={"extractor": "docling"})]


def should_retry_pdf_with_docling(units: list[SourceUnit]) -> bool:
    if not units:
        return True

    page_count = len(units)
    nonempty_units = [unit for unit in units if unit.text.strip()]
    if not nonempty_units:
        return True

    combined_text = " ".join(unit.text for unit in nonempty_units)
    avg_chars_per_page = len(combined_text) / max(page_count, 1)
    broken_word_count = len(re.findall(r"(?<=[A-Za-zА-Яа-яЁё])\s+-\s+(?=[A-Za-zА-Яа-яЁё])", combined_text))
    empty_ratio = 1 - (len(nonempty_units) / max(page_count, 1))

    if avg_chars_per_page < 120:
        return True
    if empty_ratio > 0.35:
        return True
    if broken_word_count >= max(3, page_count):
        return True
    return False


def infer_extractor(source_units: list[SourceUnit]) -> str:
    for unit in source_units:
        extractor = unit.meta.get("extractor")
        if isinstance(extractor, str) and extractor:
            return extractor
    return "default"


def combine_source_units(source_units: list[SourceUnit]) -> tuple[str, list[dict[str, int]]]:
    parts: list[str] = []
    page_spans: list[dict[str, int]] = []
    cursor = 0

    for unit in source_units:
        text = unit.text.strip()
        if not text:
            continue

        if parts:
            parts.append("\n\n")
            cursor += 2

        start = cursor
        parts.append(text)
        cursor += len(text)
        end = cursor

        page_number = unit.meta.get("page_number")
        if isinstance(page_number, int):
            page_spans.append(
                {
                    "page_number": page_number,
                    "start": start,
                    "end": end,
                }
            )

    return "".join(parts), page_spans


def resolve_page_metadata(page_spans: list[dict[str, int]], start: int, end: int) -> dict[str, object]:
    overlapping_pages = [
        span["page_number"]
        for span in page_spans
        if start < span["end"] and end > span["start"]
    ]
    if not overlapping_pages:
        return {}

    page_start = min(overlapping_pages)
    page_end = max(overlapping_pages)
    metadata: dict[str, object] = {
        "page_start": page_start,
        "page_end": page_end,
        "page_label": f"p.{page_start}" if page_start == page_end else f"pp.{page_start}-{page_end}",
    }
    if page_start == page_end:
        metadata["page_number"] = page_start
    return metadata


def extract_docx_text(path: Path) -> str:
    document = DocxDocument(str(path))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(parts)


def extract_pptx_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
    return "\n".join(parts)


def extract_xlsx_text(path: Path) -> str:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def build_metadata(path: Path, input_dir: Path, text: str) -> dict[str, object]:
    relative_path = path.relative_to(input_dir)
    parent_parts = relative_path.parts[:-1]
    domain_raw = parent_parts[0] if len(parent_parts) >= 1 else ""
    category_raw = parent_parts[1] if len(parent_parts) >= 2 else ""
    subcategory_raw = parent_parts[2] if len(parent_parts) >= 3 else ""
    source_dir = "" if relative_path.parent == Path(".") else str(relative_path.parent)

    return {
        "extension": path.suffix.lower(),
        "source_name": path.name,
        "source_name_norm": path.name.casefold(),
        "source_stem": path.stem,
        "source_stem_norm": path.stem.casefold(),
        "source_path": str(relative_path),
        "source_path_norm": str(relative_path).casefold(),
        "source_dir": source_dir,
        "source_dir_norm": source_dir.casefold(),
        "path_depth": len(parent_parts),
        "domain_raw": domain_raw,
        "domain": domain_raw.casefold(),
        "category_raw": category_raw,
        "category": category_raw.casefold(),
        "subcategory_raw": subcategory_raw,
        "subcategory": subcategory_raw.casefold(),
        "language_hint": infer_language_hint(relative_path=relative_path, text=text),
    }


def infer_language_hint(relative_path: Path, text: str) -> str:
    path_tokens = {
        token
        for token in re.split(r"[\W_]+", " ".join(relative_path.parts).casefold())
        if token
    }
    if path_tokens & RU_MARKERS:
        return "ru"
    if path_tokens & EN_MARKERS:
        return "en"

    sample = text[:4000]
    return "ru" if CYRILLIC_RE.search(sample) else "en"


def clean_pdf_text(text: str) -> str:
    cleaned = text.replace("\u00ad", "")
    cleaned = re.sub(r"(?<=\w)-\s*\n\s*(?=\w)", "", cleaned)
    # Many PDF extracts insert artificial word breaks like "усло - вия" or "рас - положено".
    # Repair them before collapsing line breaks so embeddings and exact phrase checks stay meaningful.
    cleaned = re.sub(r"(?<=[A-Za-zА-Яа-яЁё])\s+-\s+(?=[A-Za-zА-Яа-яЁё])", "", cleaned)
    cleaned = re.sub(r"(?<!\n)\s*\n\s*(?!\n)", " ", cleaned)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r" *\n{2,} *", "\n\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def clean_docling_text(text: str) -> str:
    cleaned = text.replace("\u00ad", "")
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    return cleaned.strip()


def chunk_text_with_offsets(text: str, chunk_size: int, overlap: int) -> list[ChunkSpan]:
    normalized = normalize_text(text)
    if len(normalized) <= chunk_size:
        return [ChunkSpan(content=normalized, start=0, end=len(normalized))]

    chunks: list[ChunkSpan] = []
    start = 0
    text_length = len(normalized)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        if end < text_length:
            split_at = normalized.rfind("\n\n", start, end)
            if split_at > start + chunk_size // 2:
                end = split_at
        chunk = normalized[start:end].strip()
        if chunk:
            chunk_start = normalized.find(chunk, start, end)
            chunk_end = chunk_start + len(chunk)
            chunks.append(ChunkSpan(content=chunk, start=chunk_start, end=chunk_end))
        if end >= text_length:
            break
        start = max(end - overlap, 0)

    return chunks


def normalize_text(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    return "\n".join(lines).strip()


if __name__ == "__main__":
    main()
